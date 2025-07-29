import torch
import time
from PIL import Image
import requests
import io
import json
import gradio as gr
from transformers import AutoModelForCausalLM, AutoTokenizer, TextIteratorStreamer
import os
from threading import Thread
from fastapi.middleware.cors import CORSMiddleware
from fastapi import FastAPI
from fastapi.responses import JSONResponse, StreamingResponse, Response
import pymupdf
import docx
from pptx import Presentation
from pydantic import BaseModel, Field
import uuid
from typing import Dict, List, Literal, Optional, Union
from loguru import logger


if os.environ.get("TOPS_VISIBLE_DEVICES") is not None:
    # 支持GCU算力卡
    try:
        import torch_gcu  # 导入 torch_gcu
        from torch_gcu import transfer_to_gcu  # 导入 transfer_to_gcu
        device = "gcu"
    except Exception as e:
        raise e
elif os.environ.get("CUDA_VISIBLE_DEVICES") is not None and torch.cuda.is_available():
    device = "cuda"
else:
    device = "cpu"


app = FastAPI(root_path="/openai")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

model = None


def model_load(model_path):
    global model
    if "Qwen2-VL" in model_path:
        model = Qwen2VLInference(model_path)
    else:

        model = Inference(model_path)


class Inference:
    def __init__(self, model_path):
        from transformers import AutoModelForCausalLM, AutoTokenizer, TextIteratorStreamer
        self.model = AutoModelForCausalLM.from_pretrained(
            model_path,
            torch_dtype=torch.bfloat16,
            low_cpu_mem_usage=True,
            trust_remote_code=True
        ).to(device)
        self.tokenizer = AutoTokenizer.from_pretrained(
            model_path, trust_remote_code=True)
        self.model.eval()

    def transform_messages(self, messages):
        """
        封装一个函数将
            messages = [
            {
                "role": "user",
                "content": [
                {
                    "type": "image_url",
                    "image_url": {
                        "url": ""https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-VL/assets/demo.jpeg"" 
                    },
                },
                {"type": "text", "text": "Describe this image."},
                ],
            },
            {
                "role": "user",
                "content": "Describe this image.",
            },...
            ]
        转换成
        [{"role": "user", "image": image, "content": query},{"role": "user", "content": query}]
        """
        converted_messages = []

        for message in messages:
            if "content" in message:
                # 检查 content 是否为列表
                if isinstance(message["content"], list):
                    image = None
                    query = None
                    image_url = None

                    for content in message["content"]:
                        if content.get("type") == "image_url":
                            # 提取图片 URL
                            image_url = content["image_url"]["url"]
                        elif content.get("type") == "text":
                            # 提取文本内容
                            query = content["text"]

                    # 添加转换后的消息
                    if image_url and query:
                        choice, image = mode_load(image_url)
                        converted_messages.append(
                            {"role": message["role"], "image": image, "content": query})
                    elif query:
                        converted_messages.append(
                            {"role": message["role"], "content": query})
                elif isinstance(message["content"], str):
                    # 如果 content 是字符串，直接添加
                    converted_messages.append(
                        {"role": message["role"], "content": message["content"]})

        return converted_messages

    def infer_generator(self, messages, max_length, temperature, top_p, top_k, penalty):
        messages = self.transform_messages(messages)
        # [{"role": "user", "image": image, "content": query},...] -> messages
        from transformers import TextIteratorStreamer
        input_ids = self.tokenizer.apply_chat_template(
            messages, tokenize=True, add_generation_prompt=True, return_tensors="pt", return_dict=True).to(self.model.device)
        ###############
        streamer = TextIteratorStreamer(
            self.tokenizer, timeout=60.0, skip_prompt=True, skip_special_tokens=True)

        generate_kwargs = dict(
            max_length=max_length,
            streamer=streamer,
            do_sample=True,
            top_p=top_p,
            top_k=top_k,
            temperature=temperature,
            repetition_penalty=penalty,
            eos_token_id=[151329, 151336, 151338],
        )
        gen_kwargs = {**input_ids, **generate_kwargs}

        with torch.no_grad():
            thread = Thread(target=self.model.generate, kwargs=gen_kwargs)
            thread.start()
            buffer = ""
            for new_text in streamer:
                buffer += new_text
                yield buffer


class Qwen2VLInference(Inference):
    def __init__(self, model_path):
        from transformers import Qwen2VLForConditionalGeneration, AutoTokenizer, AutoProcessor
        from qwen_vl_utils import process_vision_info
        self.model = model = Qwen2VLForConditionalGeneration.from_pretrained(
            model_path, torch_dtype="auto", device_map="auto"
        ).to(device)
        self.tokenizer = AutoProcessor.from_pretrained(
            model_path, trust_remote_code=True)
        self.model.eval()

    def transform_messages(self, messages):
        """
        封装一个函数将
            messages = [
            {
                "role": "user",
                "content": [
                {
                    "type": "image_url",
                    "image_url": {
                        "url": ""https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-VL/assets/demo.jpeg"" 
                    },
                },
                {"type": "text", "text": "Describe this image."},
                ],
            },
            {
                "role": "user",
                "content": "Describe this image.",
            },
            ]
        转换成
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "image": "https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-VL/assets/demo.jpeg",
                        },
                        {"type": "text", "text": "Describe this image."},
                    ],
                },
                {
                    "role": "user",
                    "content": "Describe this image.",
                },
                ...
            ]
        """
        for message in messages:
            if "content" in message:
                # 检查 content 是否为列表
                if isinstance(message["content"], list):
                    for content in message["content"]:
                        if content["type"] == "image_url":
                            # 更新类型和 URL
                            content["type"] = "image"
                            content["image"] = content["image_url"]["url"]
                            # 删除原来的 image_url 字段
                            del content["image_url"]
                # 如果 content 是字符串，保持不变
        return messages

    def infer_generator(self, messages, max_length, temperature, top_p, top_k, penalty):
        messages = self.transform_messages(messages)
        from transformers import TextIteratorStreamer
        from qwen_vl_utils import process_vision_info
        text = self.tokenizer.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )
        image_inputs, video_inputs = process_vision_info(messages)
        inputs_ids = self.tokenizer(
            text=[text],
            images=image_inputs,
            videos=video_inputs,
            padding=True,
            return_tensors="pt",
        ).to(self.model.device)
        ###############
        streamer = TextIteratorStreamer(
            self.tokenizer, timeout=60.0, skip_prompt=True, skip_special_tokens=True)

        generate_kwargs = dict(
            max_length=max_length,
            streamer=streamer,
            do_sample=True,
            top_p=top_p,
            top_k=top_k,
            temperature=temperature,
            repetition_penalty=penalty,

        )

        gen_kwargs = {**inputs_ids, **generate_kwargs}

        with torch.no_grad():
            thread = Thread(target=self.model.generate, kwargs=gen_kwargs)
            thread.start()
            buffer = ""
            for new_text in streamer:
                print(new_text)
                buffer += new_text
                yield buffer


def extract_text(path):
    return open(path, 'r').read()


def extract_pdf(path):
    doc = pymupdf.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_docx(path):
    doc = docx.Document(path)
    data = []
    for paragraph in doc.paragraphs:
        data.append(paragraph.text)
    content = '\n\n'.join(data)
    return content


def extract_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def mode_load(path):
    choice = ""
    file_type = path.split(".")[-1]
    print(file_type)
    if file_type in ["pdf", "txt", "py", "docx", "pptx", "json", "cpp", "md"]:
        if file_type.endswith("pdf"):
            content = extract_pdf(path)
        elif file_type.endswith("docx"):
            content = extract_docx(path)
        elif file_type.endswith("pptx"):
            content = extract_pptx(path)
        else:
            content = extract_text(path)
        choice = "doc"
        print(content[:100])
        return choice, content[:5000]
    elif file_type in ["png", "jpg", "jpeg", "bmp", "tiff", "webp"]:
        if "http" in path:
            response = requests.get(path)
            content = Image.open(io.BytesIO(response.content)).convert('RGB')
        else:
            content = Image.open(path).convert('RGB')
        choice = "image"
        return choice, content
    else:
        raise gr.Error("Oops, unsupported files.")


def stream_chat(message, history: list, temperature: float, max_length: int, top_p: float, top_k: int, penalty: float):
    print(f'message is - {message}')
    print(f'history is - {history}')
    conversation = []
    prompt_files = []
    if type(message) == list:
        conversation = message
    elif message["files"]:
        file_path = message["files"][-1]
        choice, contents = mode_load(file_path)
        if choice == "image":
            # conversation.append(
            #     {"role": "user", "image": contents, "content": message['text']})
            conversation.append(
                {"role": "user", "content":  [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": file_path
                        },
                    },
                    {"type": "text", "text": message['text']},
                ]})
        elif choice == "doc":
            format_msg = contents + "\n\n\n" + \
                "{} files uploaded.\n" + message['text']
            conversation.append({"role": "user", "content": format_msg})
    else:
        if len(history) == 0:
            # raise gr.Error("Please upload an image first.")
            contents = None
            conversation.append(
                {"role": "user", "content": message['text']})
        else:
            # image = Image.open(history[0][0][0])
            for prompt, answer in history:
                if answer is None:
                    prompt_files.append(prompt[0])
                    conversation.extend([{"role": "user", "content": ""}, {
                                        "role": "assistant", "content": ""}])
                else:
                    conversation.extend([{"role": "user", "content": prompt}, {
                                        "role": "assistant", "content": answer}])
            file_path = prompt_files[-1]
            choice, contents = mode_load(file_path)
            if choice == "image":
                # conversation.append(
                #     {"role": "user", "image": contents, "content": message['text']})
                conversation.append(
                    {"role": "user", "content":  [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": file_path
                            },
                        },
                        {"type": "text", "text": message['text']},
                    ]})
            elif choice == "doc":
                format_msg = contents + "\n\n\n" + \
                    "{} files uploaded.\n" + message['text']
                conversation.append(
                    {"role": "user", "content": format_msg})

    print(f"Conversation is -\n{conversation}")
    for text in model.infer_generator(messages=conversation, temperature=temperature, max_length=max_length, top_p=top_p, top_k=top_k, penalty=penalty):
        yield text


def demo():
    CSS = """
    h1 {
        text-align: center;
        display: block;
    }
    """

    chatbot = gr.Chatbot(label="Chatbox", height=600)
    chat_input = gr.MultimodalTextbox(
        interactive=True,
        placeholder="Enter message or upload a file one time...",
        show_label=False,

    )

    with gr.Blocks(css=CSS, theme="soft", fill_height=True) as demo:
        gr.ChatInterface(
            fn=stream_chat,
            multimodal=True,
            textbox=chat_input,
            chatbot=chatbot,
            fill_height=True,
            additional_inputs_accordion=gr.Accordion(
                label="⚙️ Parameters", open=False, render=False),
            additional_inputs=[
                gr.Slider(
                    minimum=0,
                    maximum=1,
                    step=0.1,
                    value=0.8,
                    label="Temperature",
                    render=False,
                ),
                gr.Slider(
                    minimum=1024,
                    maximum=8192,
                    step=1,
                    value=4096,
                    label="Max Length",
                    render=False,
                ),
                gr.Slider(
                    minimum=0.0,
                    maximum=1.0,
                    step=0.1,
                    value=1.0,
                    label="top_p",
                    render=False,
                ),
                gr.Slider(
                    minimum=1,
                    maximum=20,
                    step=1,
                    value=10,
                    label="top_k",
                    render=False,
                ),
                gr.Slider(
                    minimum=0.0,
                    maximum=2.0,
                    step=0.1,
                    value=1.0,
                    label="Repetition penalty",
                    render=False,
                ),
            ],
        )

    return demo


def transform_messages(messages):
    """
    封装一个函数将
        messages = [
          {
              "role": "user",
              "content": [
              {
                  "type": "image_url",
                  "image_url": {
                    "url": ""https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-VL/assets/demo.jpeg"" 
                  },
              },
              {"type": "text", "text": "Describe this image."},
              ],
              ...
          }
        ]
    转换成
    [{"role": "user", "image": image, "content": query},...]
    """
    transformed = []
    for message in messages:
        # 提取图片 URL 和文本内容
        image_url = None
        query = None

        for content in message.get("content", []):
            if content["type"] == "image_url":
                image_url = content["image_url"]["url"]
            elif content["type"] == "text":
                query = content["text"]

        # 将提取的信息添加到新的格式中
        if image_url and query:
            choice, image = mode_load(image_url)
            transformed.append(
                {"role": message["role"], "image": image, "content": query})

    return transformed


class ChatCompletesRequest(BaseModel):
    model: str = "gpt-3.5-turbo"
    messages: list
    temperature: float = 0.8
    max_tokens: int = 4096
    top_p: float = 1.0
    top_k: int = 10
    n: int = 1
    stream: bool = False
    logprobs: int = 0
    echo: bool = False
    stop: list = []
    presence_penalty: float = 1.0
    frequency_penalty: float = 1.0
    best_of: int = 1


def random_uuid() -> str:
    return str(uuid.uuid4().hex)


class UsageInfo(BaseModel):
    prompt_tokens: int = 0
    total_tokens: int = 0
    completion_tokens: int = 0
    tps: float = 0


class DeltaMessage(BaseModel):
    role: Optional[str] = None
    content: Optional[str] = None


class ChatCompletionResponseStreamChoice(BaseModel):
    index: int
    delta: DeltaMessage
    finish_reason: Optional[Literal["stop", "length"]] = None


class ChatCompletionStreamResponse(BaseModel):
    id: str = Field(default_factory=lambda: f"chatcmpl-{random_uuid()}")
    object: str = "chat.completion.chunk"
    created: int = Field(default_factory=lambda: int(time.time()))
    model: str
    choices: List[ChatCompletionResponseStreamChoice]
    usage: Optional[UsageInfo] = Field(
        default=None, description="data about request and response")


@app.post("/v1/chat/completions")
async def chat_completion(data: ChatCompletesRequest):
    logger.info(f"request: {data}")
    if data.temperature <= 0:
        data.temperature = 0.01
    if data.max_tokens <= 0:
        data.max_tokens = 1
    if data.top_p <= 0:
        data.top_p = 0.01
    if data.top_k <= 0:
        data.top_k = 10
    if data.presence_penalty <= 0:
        data.presence_penalty = 1.0
    if data.frequency_penalty <= 0:
        data.frequency_penalty = 1.0
    model = data.model

    def _gen_data(text, finish_reason=None):
        chunk = ChatCompletionStreamResponse(
            model=model,
            choices=[{
                "index": 0,
                "delta": {
                    "role": "assistant",
                    "content": text
                },
                "finish_reason": finish_reason
            }],
            usage=None
        )
        # data = chunk.json(exclude_unset=True, ensure_ascii=False)
        data = json.dumps(chunk.dict(), ensure_ascii=False)
        return f"data: {data}\n\n"

    def _gen():
        before_text = ""
        for text in stream_chat(
            message=data.messages,
            history=[],
            temperature=data.temperature,
            max_length=data.max_tokens,
            top_p=data.top_p,
            top_k=data.top_k,
            penalty=data.presence_penalty,
        ):
            yield _gen_data(text=text[len(before_text):])
            before_text = text
        yield _gen_data(text="", finish_reason="stop")
        before_text = text
    return StreamingResponse(_gen(), media_type="text/event-stream")


########################### web html ############################
app = gr.mount_gradio_app(app, demo(), path="/")


def run(port: int = 9000, model_path=""):
    import uvicorn
    model_load(model_path)
    uvicorn.run(app=app, host="0.0.0.0", port=port)


if __name__ == "__main__":
    from fire import Fire
    Fire(run)
